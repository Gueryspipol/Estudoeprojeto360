import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path

import streamlit as st
from pptx import Presentation

st.set_page_config(page_title="Estudoeprojeto360", layout="wide")
st.title("Estudoeprojeto360")
st.caption("Matriz variável, seleção na ordem desejada e geração da apresentação.")

BASE = Path(__file__).resolve().parent
SCRIPT = BASE / "projetoestudo_streamlit.py"
COMPAT = BASE / "colab_compat.py"


def localizar_modelo(novo_nome, nome_legado):
    novo = BASE / novo_nome
    legado = BASE / nome_legado
    return novo if novo.exists() else legado


MODELO_XLSX = localizar_modelo(
    "modelo_alimentador.xlsx",
    "Mezzo_Apresentação Saúde_2026_07_v1.xlsx",
)
MODELO_PPTX = localizar_modelo(
    "modelo_apresentacao.pptx",
    "Mezzo_Apresentação Saúde_2026_07_v1.pptx",
)

for caminho in (SCRIPT, COMPAT, MODELO_XLSX, MODELO_PPTX):
    if not caminho.exists():
        st.error(f"Arquivo necessário não encontrado no GitHub: {caminho.name}")
        st.stop()

nome_empresa = st.text_input("1. Nome da empresa para a capa")
matriz = st.file_uploader("2. Selecione a matriz Excel", type=["xlsx"], key="matriz")


def executar_leitura(matriz_bytes, nome_matriz):
    with tempfile.TemporaryDirectory() as temp:
        pasta = Path(temp)
        arquivo_matriz = pasta / nome_matriz
        arquivo_matriz.write_bytes(matriz_bytes)
        (pasta / SCRIPT.name).write_bytes(SCRIPT.read_bytes())
        (pasta / COMPAT.name).write_bytes(COMPAT.read_bytes())

        env = os.environ.copy()
        env["AP360_MODO_LEITURA"] = "1"
        env["AP360_UPLOAD_QUEUE"] = json.dumps([str(arquivo_matriz)], ensure_ascii=False)

        processo = subprocess.run(
            [sys.executable, SCRIPT.name],
            cwd=pasta,
            env=env,
            capture_output=True,
            text=True,
            timeout=900,
        )
        resultado = pasta / "ap360_cenarios.json"
        if processo.returncode != 0 or not resultado.exists():
            return None, processo.stderr or processo.stdout
        return json.loads(resultado.read_text(encoding="utf-8")), processo.stdout


def trocar_nome_capa(caminho_ppt, empresa):
    ppt = Presentation(caminho_ppt)
    if not ppt.slides:
        return
    slide = ppt.slides[0]
    alterado = False
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        texto = shape.text.strip()
        if texto.casefold() == "mezzo":
            for paragrafo in shape.text_frame.paragraphs:
                for run in paragrafo.runs:
                    run.text = ""
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                shape.text_frame.paragraphs[0].runs[0].text = empresa
            else:
                shape.text = empresa
            alterado = True
    if not alterado:
        # Não cria regras para empresa/operadora; apenas informa no log se a caixa da capa não existir.
        pass
    ppt.save(caminho_ppt)


if matriz:
    assinatura = f"{matriz.name}:{matriz.size}:{hash(matriz.getvalue())}"
    if st.session_state.get("assinatura") != assinatura:
        st.session_state["assinatura"] = assinatura
        st.session_state.pop("leitura", None)
        st.session_state.pop("ppt_final", None)

    if "leitura" not in st.session_state:
        with st.spinner("Lendo a matriz e identificando os cenários..."):
            leitura, log = executar_leitura(matriz.getvalue(), matriz.name)
        if leitura is None:
            st.error("O robô não conseguiu concluir a leitura da matriz.")
            with st.expander("Detalhes"):
                st.code(log)
            st.stop()
        st.session_state["leitura"] = leitura

    leitura = st.session_state["leitura"]
    atual = leitura.get("operadora_atual") or "Não identificada"
    opcoes = list(leitura.get("operadoras_disponiveis", []))

    st.info(f"Operadora atual identificada: {atual}")
    escolhidos = st.multiselect(
        "3. Escolha até quatro cenários, na ordem em que devem aparecer",
        options=opcoes,
        max_selections=4,
        placeholder="Selecione os cenários",
    )
    if escolhidos:
        st.write("Ordem da apresentação: " + " → ".join(escolhidos))

    if st.button("Gerar apresentação", type="primary", use_container_width=True):
        if not nome_empresa.strip():
            st.error("Informe o nome da empresa para a capa.")
            st.stop()
        if not escolhidos:
            st.error("Selecione pelo menos um cenário.")
            st.stop()

        # Converte os nomes escolhidos para os números que a lógica original do Colab já espera.
        numeros = [str(opcoes.index(nome) + 1) for nome in escolhidos]

        with tempfile.TemporaryDirectory() as temp:
            pasta = Path(temp)
            arquivo_matriz = pasta / matriz.name
            arquivo_matriz.write_bytes(matriz.getvalue())

            # A lógica original espera estes nomes, mas os arquivos são apenas moldes internos.
            arquivo_xlsx = pasta / "Mezzo_Apresentação Saúde_2026_07_v1.xlsx"
            arquivo_xlsx.write_bytes(MODELO_XLSX.read_bytes())
            arquivo_ppt = pasta / "Mezzo_Apresentação Saúde_2026_07_v1.pptx"
            arquivo_ppt.write_bytes(MODELO_PPTX.read_bytes())

            (pasta / SCRIPT.name).write_bytes(SCRIPT.read_bytes())
            (pasta / COMPAT.name).write_bytes(COMPAT.read_bytes())

            env = os.environ.copy()
            env.pop("AP360_MODO_LEITURA", None)
            env["AP360_UPLOAD_QUEUE"] = json.dumps(
                [str(arquivo_matriz), str(arquivo_xlsx), str(arquivo_ppt)],
                ensure_ascii=False,
            )
            env["CENARIOS_SELECIONADOS"] = ",".join(numeros)

            with st.spinner("Executando o robô original e gerando a apresentação..."):
                processo = subprocess.run(
                    [sys.executable, SCRIPT.name],
                    cwd=pasta,
                    env=env,
                    capture_output=True,
                    text=True,
                    timeout=900,
                )

            if processo.returncode != 0:
                st.error("O robô encontrou um erro durante o processamento.")
                with st.expander("Detalhes do erro"):
                    st.code(processo.stderr or processo.stdout)
                st.stop()

            saida = pasta / "Apresentador360_V2_FINAL.pptx"
            if not saida.exists():
                st.error("O robô terminou sem criar o PowerPoint final.")
                with st.expander("Registro"):
                    st.code(processo.stdout)
                st.stop()

            trocar_nome_capa(saida, nome_empresa.strip())
            st.session_state["ppt_final"] = saida.read_bytes()
            st.session_state["nome_saida"] = (
                "Apresentador360_" + nome_empresa.strip().replace(" ", "_") + ".pptx"
            )

    if "ppt_final" in st.session_state:
        st.success("Apresentação gerada com sucesso.")
        st.download_button(
            "Baixar PowerPoint",
            data=st.session_state["ppt_final"],
            file_name=st.session_state["nome_saida"],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
        st.caption("Logotipos e páginas de rede permanecem para ajuste manual.")
